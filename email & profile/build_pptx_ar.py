"""
مصلة — ملف الشركة التعريفي  |  Masalla Arabic Company Profile
نسخة كاملة عربية من التصميم الفاخر v2
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import os

# ── Palette (identical to English v2) ─────────────────────────────────────────
def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

DARK    = rgb('0C1412')
DARK2   = rgb('162220')
COPPER  = rgb('C4957A')
CREAM   = rgb('F0E8DC')
CREAM2  = rgb('EDE8E0')
MUTED   = rgb('8A9E9C')
WHITE   = rgb('FFFFFF')
PHOTO_F = rgb('CCC4BA')
PHOTO_L = rgb('B5ADA3')

IMG      = 'c:/xampp/htdocs/Masalla_Website/images'
PARTNERS = f'{IMG}/partners'

AF = 'Dubai'          # Arabic body / display font
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# ── Presentation ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
def S(): return prs.slides.add_slide(blank)

# ── Helpers ───────────────────────────────────────────────────────────────────
def box(sl, l, t, w, h, fill=None):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:    sh.fill.background()
    return sh

def overlay(sl, l, t, w, h, color='0C1412', alpha=70):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(color)
    spPr = sh._element.spPr
    srgbClr = spPr.find(f'.//{{{NS_A}}}srgbClr')
    alpha_el = etree.SubElement(srgbClr, f'{{{NS_A}}}alpha')
    alpha_el.set('val', str(int(alpha * 1000)))
    return sh

def hline(sl, l, t, w, color=None, height_pt=1.5):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Pt(height_pt))
    sh.line.fill.background()
    sh.fill.solid(); sh.fill.fore_color.rgb = color or COPPER
    return sh

def vline(sl, l, t, w_inch, h, color=None):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w_inch), Inches(h))
    sh.line.fill.background()
    sh.fill.solid(); sh.fill.fore_color.rgb = color or COPPER
    return sh

def photo(sl, l, t, w, h, tag='صورة'):
    box(sl, l, t, w, h, PHOTO_F)
    for conn in [
        sl.shapes.add_connector(1, Inches(l), Inches(t), Inches(l+w), Inches(t+h)),
        sl.shapes.add_connector(1, Inches(l+w), Inches(t), Inches(l), Inches(t+h)),
    ]:
        conn.line.color.rgb = PHOTO_L; conn.line.width = Pt(0.8)
    ar(sl, f'[ {tag} ]', l, t + h/2 - 0.18, w, 0.36,
       size=9, color=rgb('8A8078'), align=PP_ALIGN.CENTER, bold=True, font='Calibri')

def try_img(sl, path, l, t, w, h):
    if os.path.exists(path):
        try: sl.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
        except: pass

def ar(sl, text, l, t, w, h,
       size=12, bold=False, color=None, align=PP_ALIGN.RIGHT,
       italic=False, font=None):
    """RTL Arabic textbox."""
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    pPr = p._p.get_or_add_pPr(); pPr.set('rtl', '1')
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color or CREAM2
    r.font.name = font or AF
    rPr = r._r.get_or_add_rPr()
    cs = etree.SubElement(rPr, f'{{{NS_A}}}cs')
    cs.set('typeface', font or AF)
    return tb

# ── RTL mirror helper: flip x-coordinate ─────────────────────────────────────
# For RTL layout: a shape at x with width w mirrors to (13.33 - x - w)
def mx(l, w): return 13.33 - l - w


# ═══════════════════════════════════════════════════════════════════════════════
# الشريحة ١ — الغلاف
#   صورة كاملة · تراكب داكن · اسم الشركة في المنتصف
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
photo(s, 0, 0, 13.33, 7.5, 'صورة الغلاف الرئيسية')
overlay(s, 0, 0, 13.33, 7.5, '060E0C', alpha=68)

# إطار نحاسي رفيع
hline(s, 0.55, 0.5,  12.23, COPPER, 1)
hline(s, 0.55, 6.96, 12.23, COPPER, 1)

# زخرفة نجمية
ar(s, '✦', 6.1, 1.55, 1.13, 0.55, size=14, color=COPPER,
   align=PP_ALIGN.CENTER, font='Georgia')

# اسم الشركة
ar(s, 'مصلة', 1.0, 1.95, 11.33, 1.85,
   size=96, color=CREAM, align=PP_ALIGN.CENTER)
ar(s, 'MASALLA', 1.0, 3.65, 11.33, 0.9,
   size=28, color=COPPER, align=PP_ALIGN.CENTER, font='Garamond')

hline(s, 5.8, 4.72, 1.73, COPPER, 1)

ar(s, 'خ د م ا ت   ط ع ا م   ر ا ق ي ة   ·   ا ل ر ي ا ض',
   1.0, 4.95, 11.33, 0.48,
   size=10, color=rgb('C0B8B0'), align=PP_ALIGN.CENTER)

ar(s, 'hello@masalla.co   ·   www.masalla.co   ·   تأسست 2022',
   1.0, 6.55, 11.33, 0.35,
   size=9, color=rgb('7A8E8C'), align=PP_ALIGN.CENTER, font='Calibri')


# ═══════════════════════════════════════════════════════════════════════════════
# الشريحة ٢ — رسالة العلامة التجارية
#   خلفية داكنة · جملة واحدة مؤثرة · لا فوضى
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
box(s, 0, 0, 13.33, 7.5, DARK)
hline(s, 0, 0,    13.33, COPPER, 3)
hline(s, 0, 7.47, 13.33, COPPER, 3)

# حرف زخرفي ضخم في الخلفية
ar(s, 'م', 0.2, -0.5, 7.0, 8.5, size=420, color=rgb('111D1A'),
   align=PP_ALIGN.RIGHT)

# الجملة الرئيسية — RTL من اليمين
ar(s, 'لا نُقدّم وجبات\nفحسب.', 3.5, 1.3, 9.0, 3.0,
   size=62, color=CREAM, align=PP_ALIGN.RIGHT, italic=True)
ar(s, 'نصنع تجارب لا تُنسى.', 1.8, 4.05, 10.8, 1.1,
   size=62, color=COPPER, align=PP_ALIGN.RIGHT, italic=True)

hline(s, 11.43, 5.38, 1.1, COPPER, 1.5)

ar(s, 'الجودة · الاستدامة · الابتكار — في كل مستوى وكل حجم.',
   3.0, 5.62, 9.7, 0.48,
   size=13, color=MUTED, align=PP_ALIGN.RIGHT, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# الشريحة ٣ — من نحن
#   النصف الأيمن: صورة كاملة الارتفاع · النصف الأيسر: القصة + الرؤية والرسالة
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
box(s, 0, 0, 13.33, 7.5, CREAM)

# صورة — يمين (RTL mirror of left)
photo(s, 13.33-5.8, 0, 5.8, 7.5, 'صورة المطبخ / الفريق')
overlay(s, 13.33-5.8, 0, 5.8, 7.5, '0C1412', alpha=15)

# شريط نحاسي فاصل
vline(s, 13.33-5.86, 0, 0.06, 7.5, COPPER)

# منطقة النص — يسار
box(s, 0, 0, 7.47, 7.5, CREAM)

ar(s, 'من نحن', 0.5, 0.7, 6.6, 0.38,
   size=9, bold=True, color=COPPER, font='Calibri')
ar(s, 'معيار جديد\nفي ابتكار الطعام',
   0.5, 1.08, 6.7, 1.55,
   size=32, color=DARK, align=PP_ALIGN.RIGHT)
hline(s, 6.18, 2.75, 0.85, COPPER, 2)

ar(s,
   'مصلة شركة طعام مقرها الرياض تُحدث ثورة في قطاعَي B2B وB2G\n'
   'من خلال خدمات غذائية متنوعة وعالمية المستوى. من منشأتنا المركزية\n'
   'المتطورة نوفر التميز الطهوي والكفاءة اللوجستية والجودة دون تنازل\n'
   'للشركات والجهات الحكومية في أرجاء المملكة.',
   0.5, 2.98, 6.6, 1.7, size=11.5, color=MUTED, align=PP_ALIGN.RIGHT, font=AF)

# بطاقة الرؤية
box(s, 0.28, 4.9, 3.3, 2.3, DARK)
hline(s, 0.28, 4.9, 3.3, COPPER, 3)
ar(s, 'الرؤية', 0.45, 5.1, 2.8, 0.35, size=8, bold=True, color=COPPER, font='Calibri')
ar(s, 'الريادة في تقديم خدمات غذائية\nمبتكرة ومستدامة في المملكة.',
   0.45, 5.52, 2.9, 1.5, size=11, color=CREAM2, align=PP_ALIGN.RIGHT)

# بطاقة الرسالة
box(s, 3.78, 4.9, 3.3, 2.3, DARK)
hline(s, 3.78, 4.9, 3.3, COPPER, 3)
ar(s, 'الرسالة', 3.95, 5.1, 2.8, 0.35, size=8, bold=True, color=COPPER, font='Calibri')
ar(s, 'تقديم خدمات غذائية متفوقة مع\nدفع حدود الابتكار لكل عميل.',
   3.95, 5.52, 2.9, 1.5, size=11, color=CREAM2, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# الشريحة ٤ — الأرقام
#   داكنة · أربعة أرقام نحاسية ضخمة تطفو في الفضاء
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
box(s, 0, 0, 13.33, 7.5, DARK)
hline(s, 0, 0, 13.33, COPPER, 3)

ar(s, 'نظرة عامة', 7.5, 0.48, 5.2, 0.38,
   size=9, bold=True, color=COPPER, align=PP_ALIGN.RIGHT, font='Calibri')

stats = [
    ('٨+',         'خطوط الخدمة',
     'تموين مؤسسي، مطابخ سحابية، توريد هوريكا والمزيد'),
    ('B2G\n& B2B',  'القطاعات المخدومة',
     'جهات حكومية، شركات، مستشفيات، مدارس، وهوريكا'),
    ('١٠٠٪',       'صديق للبيئة',
     'مصادر مستدامة وممارسات صفرية للنفايات في جميع المنشآت'),
    ('٢٠٢٢',       'تأسست في الرياض',
     'نعمل من منشأة إنتاج مركزية متطورة تخدم كامل المملكة'),
]
# RTL: right → left (col 0 = rightmost)
col_x = [9.99, 6.77, 3.55, 0.42]
for i,(num,title,desc) in enumerate(stats):
    x = col_x[i]
    ar(s, num, x, 1.0, 3.0, 2.5,
       size=68, color=COPPER, align=PP_ALIGN.RIGHT,
       font='Garamond' if num[0].isascii() else AF)
    hline(s, x + 1.8, 3.45, 0.7, COPPER, 1)
    ar(s, title, x, 3.7, 2.9, 0.48,
       size=13, bold=True, color=CREAM2, align=PP_ALIGN.RIGHT, font=AF)
    ar(s, desc, x, 4.28, 2.9, 1.5,
       size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=AF)

for xd in [3.55, 6.77, 9.99]:
    box(s, xd, 1.0, 0.012, 5.7, rgb('2A3C39'))


# ═══════════════════════════════════════════════════════════════════════════════
# الشريحتان ٥ و٦ — خدماتنا
#   شبكة ٢×٢ — صورة تملأ ٦٥٪ العلوية · شريط داكن في الأسفل
# ═══════════════════════════════════════════════════════════════════════════════
services = [
    ('٠١', 'التموين المؤسسي',
     'وجبات مغذية للمدارس والمستشفيات والمؤسسات'),
    ('٠٢', 'خدمات الوجبات بالعقود',
     'شراكات طويلة الأمد للشركات والجهات الحكومية'),
    ('٠٣', 'تموين الفعاليات',
     'قوائم مصمّمة لفعاليات الشركات والجهات الحكومية'),
    ('٠٤', 'التموين المتنقل',
     'طعام راقٍ في أي موقع أو فعالية خارجية'),
    ('٠٥', 'توريد هوريكا',
     'منتجات طازجة ومجمّدة ومعبّأة للفنادق والمطاعم والمقاهي'),
    ('٠٦', 'علامات المطبخ السحابي',
     'مطابخ سحابية متعددة العلامات عبر تطبيقات التوصيل'),
    ('٠٧', 'إدارة الكافتيريات',
     'إدارة متكاملة — قائمة طعام، تحضير، خدمة، وتنظيف'),
    ('٠٨', 'الاستشارات',
     'سلامة غذائية، كفاءة المطبخ، وتخطيط قوائم الطعام'),
]

CARD_W = 6.38; CARD_H = 3.3; PHOTO_H = 2.25; BAR_H = CARD_H - PHOTO_H
positions = [(0.18,0.62),(6.77,0.62),(0.18,4.05),(6.77,4.05)]

for slide_idx in range(2):
    s = S()
    box(s, 0, 0, 13.33, 7.5, DARK2)
    hline(s, 0, 0, 13.33, COPPER, 3)

    part = slide_idx + 1
    ar(s, 'خدماتنا', 8.5, 0.1, 4.5, 0.38,
       size=8, bold=True, color=COPPER, align=PP_ALIGN.RIGHT, font='Calibri')
    ar(s, f'{part} / ٢', 0.1, 0.1, 1.2, 0.38,
       size=8, color=MUTED, font='Calibri', align=PP_ALIGN.LEFT)

    batch = services[slide_idx*4 : slide_idx*4+4]
    for i,(num,title,desc) in enumerate(batch):
        x, y = positions[i]

        photo(s, x, y, CARD_W, PHOTO_H, tag=f'صورة الخدمة — {title}')
        box(s, x, y + PHOTO_H, CARD_W, BAR_H, DARK)
        hline(s, x, y + PHOTO_H, CARD_W, COPPER, 2)

        # رقم الخدمة — أعلى يسار الصورة (RTL: left side)
        box(s, x, y, 0.72, 0.5, rgb('0C1412'))
        ar(s, num, x + 0.02, y + 0.02, 0.68, 0.44,
           size=14, color=COPPER, align=PP_ALIGN.CENTER)

        ar(s, title, x + 0.15, y + PHOTO_H + 0.1, CARD_W - 0.3, 0.42,
           size=13, bold=True, color=CREAM2, align=PP_ALIGN.RIGHT)
        ar(s, desc,  x + 0.15, y + PHOTO_H + 0.57, CARD_W - 0.3, 0.52,
           size=9.5, color=MUTED, align=PP_ALIGN.RIGHT, font=AF)


# ═══════════════════════════════════════════════════════════════════════════════
# الشريحة ٧ — عملاؤنا وشركاؤنا
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
box(s, 0, 0, 13.33, 7.5, CREAM)
hline(s, 0, 0, 13.33, COPPER, 3)

# شريط داكن — يمين (RTL mirror)
box(s, 13.33-4.2, 0, 4.2, 7.5, DARK)
vline(s, 13.33-4.26, 0, 0.06, 7.5, COPPER)

ar(s, 'عملاؤنا', 9.2, 1.5, 3.7, 0.85,
   size=44, color=CREAM, align=PP_ALIGN.RIGHT)
hline(s, 12.0, 3.55, 1.0, COPPER, 1.5)
ar(s, 'شركاء موثوقون في أرجاء المملكة —\nحكومي وخاص وهوريكا.',
   9.2, 3.82, 3.7, 1.1, size=11.5, color=MUTED, align=PP_ALIGN.RIGHT, font=AF)
ar(s, 'B2G  ·  B2B  ·  HORECA',
   9.2, 5.55, 3.7, 0.4, size=9, bold=True, color=COPPER, align=PP_ALIGN.RIGHT, font='Calibri')

partner_files = [
    ('MOE.png',           'وزارة التعليم'),
    ('PIF.png',           'صندوق الاستثمارات العامة'),
    ('KAFD.png',          'مركز الملك عبدالله المالي'),
    ('KAPSARC.jpg',       'كابسارك'),
    ('Elm.jpg',           'إلم'),
    ('SAB.jpg',           'بنك ساب'),
    ('MOEP.png',          'وزارة الاقتصاد'),
    ('Richard attias.jpg','ريتشارد عطية'),
]
row_y = [0.55, 3.2]
for idx,(fname,pname) in enumerate(partner_files):
    row = idx // 4; col = idx % 4
    x = 0.35 + col * 2.18
    y = row_y[row]
    box(s, x, y, 1.92, 2.35, WHITE)
    fpath = f'{PARTNERS}/{fname}'
    try_img(s, fpath, x+0.08, y+0.12, 1.76, 1.55)
    ar(s, pname, x, y+1.72, 1.92, 0.52,
       size=7.5, color=rgb('6B7E7C'), align=PP_ALIGN.CENTER, font=AF)

ar(s, '"التميّز على كل طاولة."',
   0.35, 6.55, 8.7, 0.6,
   size=13, italic=True, color=rgb('9A9290'), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# الشرائح ٨–١٠ — لحظات صُنعت باحتراف  (شريحة لكل مشروع)
# ═══════════════════════════════════════════════════════════════════════════════
projects = [
    {
        'num':    '٠١',
        'name':   'اسم المشروع',
        'type':   'نوع الفعالية',
        'client': 'العميل / الجهة',
        'date':   'الشهر والسنة',
        'guests': '٠٠٠ ضيف',
        'desc':   'وصف موجز للمشروع — المناسبة والتحدي وكيف قدّمت مصلة تجربة استثنائية مصمّمة خصيصاً لرؤية العميل.',
    },
    {
        'num':    '٠٢',
        'name':   'اسم المشروع',
        'type':   'نوع الفعالية',
        'client': 'العميل / الجهة',
        'date':   'الشهر والسنة',
        'guests': '٠٠٠ ضيف',
        'desc':   'وصف موجز للمشروع — المناسبة والتحدي وكيف قدّمت مصلة تجربة استثنائية مصمّمة خصيصاً لرؤية العميل.',
    },
    {
        'num':    '٠٣',
        'name':   'اسم المشروع',
        'type':   'نوع الفعالية',
        'client': 'العميل / الجهة',
        'date':   'الشهر والسنة',
        'guests': '٠٠٠ ضيف',
        'desc':   'وصف موجز للمشروع — المناسبة والتحدي وكيف قدّمت مصلة تجربة استثنائية مصمّمة خصيصاً لرؤية العميل.',
    },
]

for proj in projects:
    s = S()
    n = proj['num']

    photo(s, 0, 0, 13.33, 7.5, f'المشروع {n} · الصورة الرئيسية')
    overlay(s, 0, 0, 13.33, 7.5,   '060E0C', alpha=55)
    overlay(s, 0, 3.8, 13.33, 3.7, '030908', alpha=72)

    # رقم ضخم كعلامة مائية — يمين (RTL)
    ar(s, n, -0.5, -1.0, 6.0, 7.5,
       size=320, color=rgb('0E1E1B'), align=PP_ALIGN.LEFT)

    hline(s, 0, 0, 13.33, COPPER, 3)
    ar(s, 'ل ح ظ ا ت   ص ُ ن ع ت   ب ا ح ت ر ا ف',
       2.5, 0.28, 10.0, 0.38,
       size=8, bold=True, color=COPPER, align=PP_ALIGN.RIGHT, font='Calibri')
    ar(s, f'٠٣ / {n}', 0.1, 0.28, 1.5, 0.38,
       size=8, color=rgb('5A7A77'), font='Calibri', align=PP_ALIGN.LEFT)

    ar(s, proj['type'],
       0.5, 2.62, 12.0, 0.4,
       size=9, bold=True, color=COPPER, align=PP_ALIGN.RIGHT, font='Calibri')
    ar(s, proj['name'],
       0.5, 2.98, 12.5, 1.45,
       size=52, color=CREAM, align=PP_ALIGN.RIGHT, italic=True)
    hline(s, 11.78, 4.52, 1.0, COPPER, 1.5)

    # بيانات المشروع
    for mi,(mlabel,mval) in enumerate([
        ('العميل',   proj['client']),
        ('التاريخ',  proj['date']),
        ('الضيوف',   proj['guests']),
    ]):
        mx_pos = 9.2 - mi * 3.0
        ar(s, mlabel, mx_pos, 4.68, 2.7, 0.3,
           size=7.5, bold=True, color=COPPER, align=PP_ALIGN.RIGHT, font='Calibri')
        ar(s, mval, mx_pos, 4.98, 2.7, 0.35,
           size=10.5, color=CREAM2, align=PP_ALIGN.RIGHT)

    ar(s, proj['desc'],
       0.55, 5.42, 12.2, 0.9,
       size=10, italic=True, color=rgb('8AAEAA'), align=PP_ALIGN.RIGHT)

    # معرض صور صغير — أسفل يسار (RTL: left side)
    THUMB_W = 1.85; THUMB_H = 1.62; THUMB_Y = 5.72
    for ti in range(4):
        tx = 0.22 + ti * (THUMB_W + 0.08)
        photo(s, tx, THUMB_Y, THUMB_W, THUMB_H, f'معرض {ti+1}')
        overlay(s, tx, THUMB_Y, THUMB_W, THUMB_H, '060E0C', alpha=22)
        hline(s, tx, THUMB_Y, 0.35, COPPER, 2)
        vline(s, tx + THUMB_W - 0.015, THUMB_Y, 0.015, 0.35, COPPER)

    hline(s, 0, 7.47, 13.33, COPPER, 3)


# ═══════════════════════════════════════════════════════════════════════════════
# الشريحة ١١ — أختام التميّز  (الشهادات)
#   خلفية داكنة دافئة · ٦ إطارات شهادات منمّقة
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
WALL = rgb('100E0A')
box(s, 0, 0, 13.33, 7.5, WALL)
hline(s, 0, 0,    13.33, COPPER, 3)
hline(s, 0, 7.47, 13.33, COPPER, 3)

# العنوان — يمين
ar(s, 'أختام',
   7.8, 0.25, 5.1, 0.55,
   size=11, bold=True, color=COPPER, align=PP_ALIGN.RIGHT, font='Calibri')
ar(s, 'التميّز',
   6.3, 0.7, 6.5, 1.1,
   size=52, color=CREAM, align=PP_ALIGN.RIGHT, italic=True)
hline(s, 12.0, 1.9, 0.9, COPPER, 1.5)
ar(s, 'عملياتنا معتمدة وفق أعلى معايير\nسلامة الغذاء والجودة والاستدامة.',
   7.8, 2.1, 5.0, 0.9,
   size=11, color=MUTED, align=PP_ALIGN.RIGHT, font=AF)

certs = [
    ('HACCP',     'نظام سلامة\nالغذاء'),
    ('ISO 22000',  'معيار دولي\nلسلامة الغذاء'),
    ('SFDA',       'هيئة الغذاء\nوالدواء السعودية'),
    ('ISO 9001',   'نظام إدارة\nالجودة'),
    ('حلال',       'شهادة\nالحلال'),
    ('أخرى',       'أضف شهادتك\nالإضافية'),
]

FW=1.38; FH=2.55; GAP=0.08; IW=FW-GAP*2; BORDER=0.03
START_X = 0.38; SPACING = 0.22

for idx,(cname,cdesc) in enumerate(certs):
    col = idx % 3; row = idx // 3
    fx = START_X + col * (FW + SPACING)
    fy = 0.38 + row * (FH + 0.55)

    box(s, fx, fy, FW, FH, rgb('100E0A'))
    for bx,by,bw,bh in [
        (fx,           fy,           FW,     BORDER),
        (fx,           fy+FH-BORDER, FW,     BORDER),
        (fx,           fy,           BORDER, FH),
        (fx+FW-BORDER, fy,           BORDER, FH),
    ]:
        box(s, bx, by, bw, bh, COPPER)

    photo(s, fx+GAP, fy+GAP, IW, IH := IW * (FH-0.55-GAP*2)/IW, tag=cname)

    box(s, fx+GAP, fy+FH-0.6, IW, 0.52, rgb('1A1612'))
    ar(s, cname,
       fx+GAP, fy+FH-0.58, IW, 0.26,
       size=9, bold=True, color=COPPER, align=PP_ALIGN.CENTER, font='Calibri')
    ar(s, cdesc.split('\n')[1] if '\n' in cdesc else cdesc,
       fx+GAP, fy+FH-0.34, IW, 0.3,
       size=7.5, color=rgb('8A8078'), align=PP_ALIGN.CENTER, font=AF)

ar(s, 'الوثائق الكاملة للاعتمادات متاحة عند الطلب.',
   0.42, 6.98, 12.5, 0.32,
   size=8.5, italic=True, color=rgb('4A5A58'), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# الشريحة ١٢ — تواصل معنا
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
photo(s, 0, 0, 13.33, 7.5, 'صورة الختام / الطعام')
overlay(s, 0, 0, 13.33, 7.5, '060E0C', alpha=78)

hline(s, 0.6, 0.55, 12.13, COPPER, 1)
hline(s, 0.6, 6.9,  12.13, COPPER, 1)

ar(s, 'تواصل معنا',
   1.0, 1.05, 11.33, 0.45,
   size=10, bold=True, color=COPPER, align=PP_ALIGN.CENTER, font='Calibri')

ar(s, 'لنبنِ معاً\nشيئاً استثنائياً',
   1.0, 1.55, 11.33, 2.2,
   size=58, color=CREAM, align=PP_ALIGN.CENTER, italic=True)

hline(s, 6.0, 3.88, 1.33, COPPER, 1.5)

ar(s, 'تواصل معنا لمناقشة كيف يمكن لمصلة خدمة مؤسستكم.',
   1.0, 4.1, 11.33, 0.48,
   size=13, italic=True, color=rgb('B0C0BE'), align=PP_ALIGN.CENTER)

contacts = [
    ('البريد',    'hello@masalla.co'),
    ('الموقع',    'www.masalla.co'),
    ('الموقع الجغرافي', 'الرياض، المملكة العربية السعودية'),
]
for i,(label_t,val) in enumerate(contacts):
    x = 1.8 + i * 3.5
    box(s, x, 5.0, 3.0, 1.2, rgb('0C1412'))
    hline(s, x, 5.0, 3.0, COPPER, 2)
    ar(s, label_t, x, 5.1, 3.0, 0.35,
       size=8, bold=True, color=COPPER, align=PP_ALIGN.CENTER, font='Calibri')
    ar(s, val, x, 5.52, 3.0, 0.55,
       size=11, color=CREAM2, align=PP_ALIGN.CENTER,
       font='Calibri' if val[0].isascii() else AF)

ar(s, 'مصلة  ·  إعادة تعريف خدمات الطعام  ·  الرياض ٢٠٢٢',
   1.0, 7.05, 11.33, 0.32,
   size=8, color=rgb('4A6A67'), align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = 'c:/xampp/htdocs/Masalla_Website/Masalla_Company_Profile_AR.pptx'
prs.save(out)
print(f'Saved: {out}')
