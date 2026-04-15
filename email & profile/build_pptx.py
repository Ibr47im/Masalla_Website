"""
Masalla Company Profile — Luxury Editorial Build
Layout philosophy:
  · Photos are heroes — text overlays or flanks them
  · Minimal copy — only what's essential
  · Dark + copper = luxury warmth
  · Generous breathing room — never crowded
  · Each slide has one visual idea, not many
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import os

# ── Palette ───────────────────────────────────────────────────────────────────
def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

DARK    = rgb('0C1412')   # forest black-green — main dark bg
DARK2   = rgb('162220')   # slightly lighter dark panel
COPPER  = rgb('C4957A')   # brand copper / gold
CREAM   = rgb('F0E8DC')   # warm cream — light slides
CREAM2  = rgb('EDE8E0')   # body text on dark
MUTED   = rgb('8A9E9C')   # secondary text
WHITE   = rgb('FFFFFF')
PHOTO_F = rgb('CCC4BA')   # photo placeholder fill
PHOTO_L = rgb('B5ADA3')   # photo placeholder cross lines

IMG      = 'c:/xampp/htdocs/Masalla_Website/images'
PARTNERS = f'{IMG}/partners'

# ── Presentation ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def S(): return prs.slides.add_slide(blank)

# ── Core shapes ───────────────────────────────────────────────────────────────
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def box(sl, l, t, w, h, fill=None):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    return sh

def overlay(sl, l, t, w, h, color='0C1412', alpha=70):
    """Semi-transparent rectangle for photo overlays."""
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(color)  # sets solidFill in XML
    # inject alpha into the srgbClr element
    spPr = sh._element.spPr
    srgbClr = spPr.find(f'.//{{{NS_A}}}srgbClr')
    alpha_el = etree.SubElement(srgbClr, f'{{{NS_A}}}alpha')
    alpha_el.set('val', str(int(alpha * 1000)))
    return sh

def line(sl, l, t, w, color=None, height_pt=1.5):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Pt(height_pt))
    sh.line.fill.background()
    sh.fill.solid(); sh.fill.fore_color.rgb = color or COPPER
    return sh

def txt(sl, text, l, t, w, h,
        size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
        italic=False, font='Garamond'):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = color or CREAM2
    return tb

def photo(sl, l, t, w, h, tag='PHOTO'):
    """Placeholder: grey box + X cross + label."""
    box(sl, l, t, w, h, PHOTO_F)
    for conn in [
        sl.shapes.add_connector(1, Inches(l), Inches(t), Inches(l+w), Inches(t+h)),
        sl.shapes.add_connector(1, Inches(l+w), Inches(t), Inches(l), Inches(t+h)),
    ]:
        conn.line.color.rgb = PHOTO_L; conn.line.width = Pt(0.8)
    txt(sl, f'[ {tag} ]', l, t + h/2 - 0.18, w, 0.36,
        size=9, color=rgb('8A8078'), align=PP_ALIGN.CENTER,
        bold=True, font='Calibri')

def try_img(sl, path, l, t, w, h):
    if os.path.exists(path):
        try: sl.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
        except: pass

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
#   Full-bleed photo · dark overlay · company name centred
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
photo(s, 0, 0, 13.33, 7.5, 'HERO / BRAND PHOTO')
overlay(s, 0, 0, 13.33, 7.5, '060E0C', alpha=68)

# thin copper frame lines
line(s, 0.55, 0.5,  12.23, COPPER, 1)
line(s, 0.55, 6.96, 12.23, COPPER, 1)
box(s, 0.55, 0.5,  Pt(1.5)/914400*13.33, 6.46, COPPER)   # left vertical
box(s, Inches(12.78)-Pt(0.75), Inches(0.5), Pt(1.5)/914400*13.33*0, 0, COPPER)

# copper dot ornament
txt(s, '✦', 6.1, 1.55, 1.13, 0.55, size=14, color=COPPER,
    align=PP_ALIGN.CENTER, font='Georgia')

# Company name
txt(s, 'MASALLA', 1.0, 2.05, 11.33, 1.7,
    size=88, bold=False, color=CREAM, align=PP_ALIGN.CENTER, font='Garamond')
txt(s, 'مصلة', 1.0, 3.6, 11.33, 1.0,
    size=36, color=COPPER, align=PP_ALIGN.CENTER, font='Times New Roman')

line(s, 5.8, 4.72, 1.73, COPPER, 1)

txt(s, 'B E S P O K E   F O O D   S E R V I C E S   ·   R I Y A D H',
    1.0, 4.95, 11.33, 0.45,
    size=10, color=rgb('C0B8B0'), align=PP_ALIGN.CENTER, font='Garamond')

txt(s, 'hello@masalla.co   ·   www.masalla.co   ·   Est. 2022',
    1.0, 6.55, 11.33, 0.35,
    size=9, color=rgb('7A8E8C'), align=PP_ALIGN.CENTER, font='Calibri')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — BRAND STATEMENT
#   Dark · one powerful sentence · no clutter
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
box(s, 0, 0, 13.33, 7.5, DARK)
line(s, 0, 0, 13.33, COPPER, 3)
line(s, 0, 7.47, 13.33, COPPER, 3)

# large ambient letter
txt(s, 'M', 7.2, 0.0, 6.0, 7.5, size=420, color=rgb('111D1A'),
    align=PP_ALIGN.LEFT, font='Garamond', bold=False)

txt(s, 'We don\'t just\ncater events.', 0.8, 1.3, 9.0, 3.0,
    size=62, color=CREAM, align=PP_ALIGN.LEFT, font='Garamond', italic=True)
txt(s, 'We craft experiences.', 0.8, 4.05, 9.5, 1.1,
    size=62, color=COPPER, align=PP_ALIGN.LEFT, font='Garamond', italic=True)

line(s, 0.8, 5.38, 1.1, COPPER, 1.5)

txt(s, 'Quality · Sustainability · Innovation — across every scale.',
    0.8, 5.62, 9.0, 0.45,
    size=13, color=MUTED, align=PP_ALIGN.LEFT, font='Garamond', italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ABOUT
#   Left half: full-height photo · Right half: story + vision/mission
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
box(s, 0, 0, 13.33, 7.5, CREAM)

# left photo — full height
photo(s, 0, 0, 5.8, 7.5, 'KITCHEN / TEAM PHOTO')
overlay(s, 0, 0, 5.8, 7.5, '0C1412', alpha=15)  # very subtle tint

# copper accent strip between photo and text
box(s, 5.8, 0, 0.06, 7.5, COPPER)

# right text area
box(s, 5.86, 0, 7.47, 7.5, CREAM)

txt(s, 'WHO WE ARE', 6.3, 0.7, 6.6, 0.38,
    size=9, bold=True, color=COPPER, font='Calibri')
txt(s, 'A New Standard\nin Food Innovation',
    6.2, 1.08, 6.7, 1.55,
    size=32, color=DARK, align=PP_ALIGN.LEFT, font='Garamond')
line(s, 6.3, 2.75, 0.85, COPPER, 2)

txt(s,
    'Masalla is a Riyadh-based food company revolutionising B2B and B2G\n'
    'food services. From our state-of-the-art central facility, we deliver\n'
    'culinary excellence, efficient logistics, and zero-compromise quality\n'
    'to businesses and government entities across the Kingdom.',
    6.3, 2.98, 6.6, 1.7, size=11.5, color=MUTED, font='Calibri')

# Vision card
box(s, 6.2, 4.9, 3.3, 2.3, DARK)
line(s, 6.2, 4.9, 3.3, COPPER, 3)
txt(s, 'VISION', 6.45, 5.1, 2.8, 0.35, size=8, bold=True, color=COPPER, font='Calibri')
txt(s, 'The leading provider of innovative,\nsustainable food services in the Kingdom.',
    6.45, 5.52, 2.9, 1.5, size=11, color=CREAM2, font='Garamond')

# Mission card
box(s, 9.7, 4.9, 3.3, 2.3, DARK)
line(s, 9.7, 4.9, 3.3, COPPER, 3)
txt(s, 'MISSION', 9.95, 5.1, 2.8, 0.35, size=8, bold=True, color=COPPER, font='Calibri')
txt(s, 'Delivering superior food services while\npushing culinary innovation for every client.',
    9.95, 5.52, 2.9, 1.5, size=11, color=CREAM2, font='Garamond')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — NUMBERS
#   Dark · four large copper stats floating in space
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
box(s, 0, 0, 13.33, 7.5, DARK)
line(s, 0, 0, 13.33, COPPER, 3)

txt(s, 'AT A GLANCE', 0.6, 0.48, 5, 0.38,
    size=9, bold=True, color=COPPER, font='Calibri')

stats = [
    ('8+',        'Service Lines',
     'Institutional catering, cloud kitchens, HORECA supply and more'),
    ('B2G\n& B2B', 'Sectors Served',
     'Government, corporates, hospitals, schools, and HORECA'),
    ('100%',      'Eco-Friendly',
     'Sustainable sourcing and zero-waste kitchen practices across all facilities'),
    ('2022',      'Est. Riyadh',
     'Operating from a state-of-the-art central production facility'),
]
for i,(num,title,desc) in enumerate(stats):
    x = 0.42 + i * 3.22
    # large number
    txt(s, num, x, 1.0, 3.0, 2.5,
        size=72, color=COPPER, align=PP_ALIGN.LEFT, font='Garamond')
    line(s, x, 3.45, 0.7, COPPER, 1)
    txt(s, title, x, 3.7, 2.9, 0.48,
        size=13, bold=True, color=CREAM2, font='Calibri')
    txt(s, desc,  x, 4.28, 2.9, 1.5,
        size=10, color=MUTED, font='Calibri')

# subtle vertical dividers
for xd in [3.55, 6.77, 9.99]:
    box(s, xd, 1.0, 0.012, 5.7, rgb('2A3C39'))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES 5 & 6 — SERVICES
#   2×2 photo cards — photo fills top 65%, dark bar at base with service name
# ═══════════════════════════════════════════════════════════════════════════════
services = [
    ('01', 'Institutional Catering',
     'Nutritious meals for schools, hospitals & institutions'),
    ('02', 'Contract Meal Services',
     'Long-term partnerships for businesses & government entities'),
    ('03', 'Events Catering',
     'Tailored menus for corporate & government events'),
    ('04', 'Mobile Catering',
     'Top-quality food delivered to any site or outdoor location'),
    ('05', 'HORECA Supply',
     'Fresh, frozen & packaged F&B for hotels, restaurants & cafes'),
    ('06', 'Cloud Kitchen Brands',
     'Multi-brand cloud kitchens delivering B2C via aggregator apps'),
    ('07', 'Canteen Operations',
     'Full canteen management — menu, prep, service & cleanup'),
    ('08', 'Consultancy',
     'Food safety, kitchen efficiency & menu planning expertise'),
]

CARD_W = 6.38
CARD_H = 3.3
PHOTO_H = 2.25   # top portion = photo
BAR_H   = CARD_H - PHOTO_H

positions = [
    (0.18, 0.62), (6.77, 0.62),
    (0.18, 4.05), (6.77, 4.05),
]

for slide_idx in range(2):
    s = S()
    box(s, 0, 0, 13.33, 7.5, DARK2)
    line(s, 0, 0, 13.33, COPPER, 3)

    part = slide_idx + 1
    txt(s, 'OUR SERVICES', 0.22, 0.1, 4, 0.38,
        size=8, bold=True, color=COPPER, font='Calibri')
    txt(s, f'{part} / 2', 12.6, 0.1, 0.6, 0.38,
        size=8, color=MUTED, font='Calibri', align=PP_ALIGN.RIGHT)

    batch = services[slide_idx*4 : slide_idx*4+4]

    for i,(num,title,desc) in enumerate(batch):
        x, y = positions[i]

        # photo area
        photo(s, x, y, CARD_W, PHOTO_H,
              tag=f'SERVICE PHOTO — {title.upper()}')

        # dark info bar
        box(s, x, y + PHOTO_H, CARD_W, BAR_H, DARK)
        line(s, x, y + PHOTO_H, CARD_W, COPPER, 2)

        # service number — copper, top-right of photo
        box(s, x + CARD_W - 0.78, y, 0.72, 0.5, rgb('0C1412'))
        txt(s, num, x + CARD_W - 0.76, y + 0.02, 0.7, 0.44,
            size=14, color=COPPER, align=PP_ALIGN.CENTER, font='Garamond')

        # title + desc in info bar
        txt(s, title, x + 0.22, y + PHOTO_H + 0.1, CARD_W - 0.3, 0.42,
            size=13, bold=True, color=CREAM2, font='Garamond')
        txt(s, desc,  x + 0.22, y + PHOTO_H + 0.57, CARD_W - 0.3, 0.52,
            size=9.5, color=MUTED, font='Calibri')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CLIENTS & PARTNERS
#   Cream · logo grid · clean and airy
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
box(s, 0, 0, 13.33, 7.5, CREAM)
line(s, 0, 0, 13.33, COPPER, 3)

# dark left strip with heading
box(s, 0, 0, 4.2, 7.5, DARK)
box(s, 4.2, 0, 0.06, 7.5, COPPER)

txt(s, 'OUR\nCLIENTELE', 0.4, 1.5, 3.5, 1.8,
    size=44, color=CREAM, align=PP_ALIGN.LEFT, font='Garamond')
line(s, 0.4, 3.55, 1.0, COPPER, 1.5)
txt(s, 'Trusted partners across\nSaudi Arabia — government,\ncorporate, and HORECA.',
    0.4, 3.82, 3.4, 1.4, size=11.5, color=MUTED, font='Calibri')
txt(s, 'B2G  ·  B2B  ·  HORECA',
    0.4, 5.55, 3.4, 0.4, size=9, bold=True, color=COPPER, font='Calibri')

# logo grid
partner_files = [
    ('MOE.png',          'Ministry of Education'),
    ('PIF.png',          'Public Investment Fund'),
    ('KAFD.png',         'KAFD'),
    ('KAPSARC.jpg',      'KAPSARC'),
    ('Elm.jpg',          'Elm'),
    ('SAB.jpg',          'SAB Bank'),
    ('MOEP.png',         'Ministry of Economy'),
    ('Richard attias.jpg','Richard Attias & Assoc.'),
]
row_y = [0.55, 3.2]
for idx,(fname,pname) in enumerate(partner_files):
    row = idx // 4; col = idx % 4
    x = 4.55 + col * 2.18
    y = row_y[row]
    box(s, x, y, 1.92, 2.35, WHITE)
    fpath = f'{PARTNERS}/{fname}'
    try_img(s, fpath, x+0.08, y+0.12, 1.76, 1.55)
    txt(s, pname, x, y+1.72, 1.92, 0.52,
        size=7.5, color=rgb('6B7E7C'), align=PP_ALIGN.CENTER, font='Calibri')

# bottom quote
txt(s, '"Excellence served at every table."',
    4.55, 6.55, 8.4, 0.6,
    size=13, italic=True, color=rgb('9A9290'), align=PP_ALIGN.CENTER, font='Garamond')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES 8–10 — CRAFTED MOMENTS  (one slide per project)
#
#  Layout per slide:
#   · Full-bleed hero photo as backdrop with heavy dark overlay
#   · Project number as a massive faded watermark
#   · Project name large over the photo, left-anchored
#   · 4 gallery thumbnails in a horizontal strip at the bottom
#   · Thin copper details throughout
# ═══════════════════════════════════════════════════════════════════════════════

projects = [
    {
        'num':   '01',
        'name':  'Project Name',
        'type':  'EVENT TYPE',
        'client':'Client / Organisation',
        'date':  'Month Year',
        'guests':'000 Guests',
        'desc':  'A brief description of the project — the occasion, the challenge, and how Masalla delivered an exceptional experience tailored to the client\'s vision.',
    },
    {
        'num':   '02',
        'name':  'Project Name',
        'type':  'EVENT TYPE',
        'client':'Client / Organisation',
        'date':  'Month Year',
        'guests':'000 Guests',
        'desc':  'A brief description of the project — the occasion, the challenge, and how Masalla delivered an exceptional experience tailored to the client\'s vision.',
    },
    {
        'num':   '03',
        'name':  'Project Name',
        'type':  'EVENT TYPE',
        'client':'Client / Organisation',
        'date':  'Month Year',
        'guests':'000 Guests',
        'desc':  'A brief description of the project — the occasion, the challenge, and how Masalla delivered an exceptional experience tailored to the client\'s vision.',
    },
]

for proj in projects:
    s = S()
    n = proj['num']

    # ── Full-bleed hero photo ─────────────────────────────────────────────────
    photo(s, 0, 0, 13.33, 7.5, f'PROJECT {n} · HERO PHOTO')

    # Two-layer overlay: full dark + extra heavy at bottom for text legibility
    overlay(s, 0, 0, 13.33, 7.5,  '060E0C', alpha=55)
    overlay(s, 0, 3.8, 13.33, 3.7, '030908', alpha=72)

    # ── Giant watermark number ────────────────────────────────────────────────
    txt(s, n, 8.5, -1.0, 5.5, 7.5,
        size=320, color=rgb('0E1E1B'), align=PP_ALIGN.LEFT, font='Garamond')

    # ── Section label + copper top bar ───────────────────────────────────────
    line(s, 0, 0, 13.33, COPPER, 3)
    txt(s, 'C R A F T E D   M O M E N T S',
        0.55, 0.28, 6, 0.38,
        size=8, bold=True, color=COPPER, font='Calibri')
    txt(s, f'{n} / 03', 12.55, 0.28, 0.7, 0.38,
        size=8, color=rgb('5A7A77'), font='Calibri', align=PP_ALIGN.RIGHT)

    # ── Project name ──────────────────────────────────────────────────────────
    txt(s, proj['type'],
        0.55, 2.62, 8, 0.4,
        size=9, bold=True, color=COPPER, font='Calibri')
    txt(s, proj['name'],
        0.5, 2.98, 9.5, 1.45,
        size=52, color=CREAM, align=PP_ALIGN.LEFT, font='Garamond', italic=True)
    line(s, 0.55, 4.52, 1.0, COPPER, 1.5)

    # ── Project meta row ──────────────────────────────────────────────────────
    for mi, (mlabel, mval) in enumerate([
        ('CLIENT',  proj['client']),
        ('DATE',    proj['date']),
        ('GUESTS',  proj['guests']),
    ]):
        mx = 0.55 + mi * 3.0
        txt(s, mlabel, mx, 4.68, 2.7, 0.3,
            size=7.5, bold=True, color=COPPER, font='Calibri')
        txt(s, mval,   mx, 4.98, 2.7, 0.35,
            size=10.5, color=CREAM2, font='Garamond')

    # ── Description ───────────────────────────────────────────────────────────
    txt(s, proj['desc'],
        0.55, 5.42, 8.5, 0.9,
        size=10, italic=True, color=rgb('8AAEAA'), font='Garamond')

    # ── Gallery strip — 4 thumbnails at the bottom-right ─────────────────────
    THUMB_W = 1.85
    THUMB_H = 1.62
    THUMB_Y = 5.72
    for ti in range(4):
        tx = 9.18 + ti * (THUMB_W + 0.08)
        # keep last thumb within slide
        if tx + THUMB_W > 13.2:
            THUMB_W = 13.18 - tx
        photo(s, tx, THUMB_Y, THUMB_W, THUMB_H, f'GALLERY {ti+1}')
        overlay(s, tx, THUMB_Y, THUMB_W, THUMB_H, '060E0C', alpha=22)
        # copper corner accent
        line(s, tx, THUMB_Y, 0.35, COPPER, 2)
        box(s, tx, THUMB_Y, 0.015, 0.35, COPPER)

    line(s, 0, 7.47, 13.33, COPPER, 3)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SEALS OF EXCELLENCE  (Certificates)
#
#  Layout:
#   · Deep warm-dark background — like a gallery wall
#   · Section title typeset large and dramatically
#   · 6 framed certificate slots — portrait orientation, copper double-border
#   · Each frame: thin outer copper border + inner gap + white interior
#   · Certificate name beneath each frame in small elegant caps
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
WALL = rgb('100E0A')   # very dark warm near-black
box(s, 0, 0, 13.33, 7.5, WALL)

# copper top + bottom rules
line(s, 0, 0,    13.33, COPPER, 3)
line(s, 0, 7.47, 13.33, COPPER, 3)

# ── Headline ─────────────────────────────────────────────────────────────────
txt(s, 'SEALS OF',
    0.55, 0.25, 5, 0.55,
    size=11, bold=True, color=COPPER, font='Calibri')
txt(s, 'Excellence',
    0.42, 0.7, 6.5, 1.1,
    size=52, color=CREAM, align=PP_ALIGN.LEFT, font='Garamond', italic=True)
line(s, 0.55, 1.9, 0.9, COPPER, 1.5)
txt(s, 'Our operations are certified to the highest\nstandards of food safety, quality & sustainability.',
    0.55, 2.1, 4.5, 0.85,
    size=11, color=MUTED, font='Calibri')

# ── 6 framed certificate slots ───────────────────────────────────────────────
certs = [
    ('HACCP',       'Food Safety\nManagement'),
    ('ISO 22000',   'Int\'l Food\nSafety Standard'),
    ('SFDA',        'Saudi Food &\nDrug Authority'),
    ('ISO 9001',    'Quality\nManagement'),
    ('HALAL',       'Halal\nCertification'),
    ('OTHER',       'Add your\ncertificate'),
]

FW = 1.38   # frame outer width
FH = 2.55   # frame outer height
GAP = 0.08  # gap between outer and inner border
IW = FW - GAP*2
IH = FH - GAP*2

# arrange 6 frames: 3 top row, 3 bottom row, centered in right 8" of slide
START_X = 5.1
SPACING = 0.22

for idx, (cname, cdesc) in enumerate(certs):
    col = idx % 3
    row = idx // 3
    fx = START_X + col * (FW + SPACING)
    fy = 0.38 + row * (FH + 0.55)

    # outer copper frame
    box(s, fx, fy, FW, FH, rgb('100E0A'))
    # copper border — draw as 4 thin rects (top, bottom, left, right)
    BORDER = 0.03
    for bx, by, bw, bh in [
        (fx,           fy,           FW,    BORDER),  # top
        (fx,           fy+FH-BORDER, FW,    BORDER),  # bottom
        (fx,           fy,           BORDER,FH),      # left
        (fx+FW-BORDER, fy,           BORDER,FH),      # right
    ]:
        box(s, bx, by, bw, bh, COPPER)

    # inner white area for certificate image
    photo(s, fx+GAP, fy+GAP, IW, IH - 0.55, f'{cname}')

    # name beneath photo inside frame
    box(s, fx+GAP, fy+FH-0.6, IW, 0.52, rgb('1A1612'))
    txt(s, cname,
        fx+GAP, fy+FH-0.58, IW, 0.26,
        size=9, bold=True, color=COPPER, align=PP_ALIGN.CENTER, font='Calibri')
    txt(s, cdesc,
        fx+GAP, fy+FH-0.34, IW, 0.3,
        size=7.5, color=rgb('8A8078'), align=PP_ALIGN.CENTER, font='Calibri')

# ── Subtle tagline ────────────────────────────────────────────────────────────
txt(s, 'Full certification documentation available upon request.',
    0.42, 6.98, 12.5, 0.32,
    size=8.5, italic=True, color=rgb('4A5A58'), align=PP_ALIGN.CENTER, font='Garamond')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CTA / CONTACT
#   Full-bleed photo · strong overlay · minimal centred contact
# ═══════════════════════════════════════════════════════════════════════════════
s = S()
photo(s, 0, 0, 13.33, 7.5, 'CLOSING / FOOD PHOTO')
overlay(s, 0, 0, 13.33, 7.5, '060E0C', alpha=78)

# frame lines
line(s, 0.6, 0.55, 12.13, COPPER, 1)
line(s, 0.6, 6.9, 12.13, COPPER, 1)

txt(s, 'GET IN TOUCH',
    1.0, 1.05, 11.33, 0.45,
    size=10, bold=True, color=COPPER, align=PP_ALIGN.CENTER, font='Calibri')

txt(s, "Let's Build\nSomething Exceptional",
    1.0, 1.55, 11.33, 2.2,
    size=58, color=CREAM, align=PP_ALIGN.CENTER, font='Garamond', italic=True)

line(s, 6.0, 3.88, 1.33, COPPER, 1.5)

txt(s, 'Reach out to discuss how Masalla can serve your organisation.',
    1.0, 4.1, 11.33, 0.48,
    size=13, italic=True, color=rgb('B0C0BE'), align=PP_ALIGN.CENTER, font='Garamond')

# contact row
contacts = [
    ('EMAIL',    'hello@masalla.co'),
    ('WEBSITE',  'www.masalla.co'),
    ('LOCATION', 'Riyadh, Saudi Arabia'),
]
for i,(label_t, val) in enumerate(contacts):
    x = 1.8 + i * 3.5
    box(s, x, 5.0, 3.0, 1.2, rgb('0C1412'))
    line(s, x, 5.0, 3.0, COPPER, 2)
    txt(s, label_t, x, 5.1, 3.0, 0.35,
        size=8, bold=True, color=COPPER, align=PP_ALIGN.CENTER, font='Calibri')
    txt(s, val, x, 5.52, 3.0, 0.55,
        size=11, color=CREAM2, align=PP_ALIGN.CENTER, font='Garamond')

txt(s, 'Masalla  ·  Redefining Food Services  ·  Riyadh 2022',
    1.0, 7.05, 11.33, 0.32,
    size=8, color=rgb('4A6A67'), align=PP_ALIGN.CENTER, font='Calibri')


# ── Save ──────────────────────────────────────────────────────────────────────
out = 'c:/xampp/htdocs/Masalla_Website/Masalla_Company_Profile_v2.pptx'
prs.save(out)
print(f'Saved: {out}')
